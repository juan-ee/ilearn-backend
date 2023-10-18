from pptx import Presentation
from pptx.util import Inches
import copy


def fill_data(data, shape):
    paragraphs = shape.text_frame.paragraphs
    original_font = copy.deepcopy(paragraphs[1].runs[0].font)

    for index, value in data.items():
        paragraphs[index].text = value
        paragraphs[index].runs[0].font.size = original_font.size


def fill_bullet_list_data(text, shape):
    paragraphs = shape.text_frame.paragraphs
    original_font = copy.deepcopy(paragraphs[1].runs[0].font)

    paragraphs[1].text = text
    for run in paragraphs[1].runs:
        run.font.size = original_font.size


def fill_env_data(data, shapes):
    for shape in shapes:
        if shape.name == 'General':
            fill_bullet_list_data(data['general'], shape)


if __name__ == '__main__':
    general_info = {
        1: 'Henkel',
        3: 'Chemical',
        5: 'Duesseldorf',
        7: '519500'
    }

    ratings_info = {
        1: 'Platinum Medal',
        4: 'A-',
        7: '18.5 ESG Risk Rating (Low Risk)',
        10: 'AAA',
        13: 'None'
    }

    risks = '\n- '.join([
        '- Physical risks associated with extreme weather',
        'Main risk here is the kernel oil production and mainly the risks of a low crop yield because of El Nino',
        'Transition risks associated with the transition to low-emission economy and society. This risk mainly steam from a potential increase in costs associated with CO2 emissions '
    ])

    opportunities = '\n- '.join([
        '- Relaunch shampoo and conditioner brand with more recycled plastic',
        'Has published a 2030+ sustainability ambition framework with where they currently work on improving and have ambitions of improving'
        'Provide a comprehensive sustainability profile for all products by 2025',
        'Reducing material in packaging',
        'Invests in fund for sustainable packaging',
        'Collaborates with “Plastic bank” to reduce plastic waste by developing recycling infrastructure in developing countries',
        'Actively wants to contribute to thriving communities',
        'Expand community education and volunteering programs',
        'Try to engage and empower all employees to be more sustainable',
        'Supply heat generated from their production to external users for heating in office buildings',
    ])

    env_data = {
        'general': '\n- '.join([
            '- Increased traceability rate by 5% points for palm-based raw materials',
            'Company has renewed its renewable energy contract',
            'Further optimize transportation and logistic chain:',
            'by analyzing the carbon footprint of logistics and assess how achieve further CO2 reductions in 18 countries Henkel operates in',
            'Increase ocean freight and decrease air freight',
            'invest in alternative drive trains. Ex. Battery and hydrogen powered ones',
            'Expand the digital tools to increase efficiency of logistics'
        ])
    }

    presentation = Presentation('template.pptx')
    slide = presentation.slides[0]

    for shape in slide.shapes:
        if shape.name == 'GeneralInfo':
            fill_data(general_info, shape)
        if shape.name == 'Ratings':
            fill_data(ratings_info, shape)
        if shape.name == 'Risks':
            fill_bullet_list_data(risks, shape)
        if shape.name == 'Opportunities':
            fill_bullet_list_data(opportunities, shape)
        if shape.name == 'Enviromental':
            fill_env_data(env_data, shape.shapes)

    
    img_path = "C:/Users/crugi/Pictures/Wallpapers/Wooper.png"
    pic = slide.shapes.add_picture(img_path, left=Inches(0.15), top=Inches(0.23), width=Inches(1), height=Inches(0.7))
    
    presentation.save('output.pptx')
