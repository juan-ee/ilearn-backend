from pptx import Presentation
from pptx.util import Inches
import copy
from ai import *


def fill_data(data, shape):
    paragraphs = shape.text_frame.paragraphs
    original_font = copy.deepcopy(paragraphs[1].runs[0].font)

    for index, value in data.items():
        paragraphs[index].text = value
        paragraphs[index].runs[0].font.size = original_font.size


def fill_bullet_list_data(text, shape):
    paragraphs = shape.text_frame.paragraphs
    original_font = copy.deepcopy(paragraphs[1].runs[0].font)

    paragraphs[1].text = text
    for run in paragraphs[1].runs:
        run.font.size = original_font.size


def fill_env_data(data, shapes):
    for shape in shapes:
        if shape.name == 'EmissionManagement':
            fill_bullet_list_data(data['emission_management'], shape)
        if shape.name == 'ResourcesManagement':
            fill_bullet_list_data(data['resources_management'], shape)
        if shape.name == 'WasteManagement':
            fill_bullet_list_data(data['waste_management'], shape)


def fill_logo(slide, logo_path):
    slide.shapes.add_picture(logo_path, left=Inches(0.15), top=Inches(0.23), width=Inches(1), height=Inches(0.7))


def save_pptx(company_name, logo_path, ratings, pptx_path):
    general_info = {
        1: company_name,
        3: get_industry_auto(company_name.lower()),
        5: get_location(company_name.lower()),
        7: get_head_count()
    }

    ratings_info = {
        1: ratings['ecovadis'],
        4: ratings['cdp'],
        7: ratings['sustainalytics'],
        10: ratings['msci'],
        13: ratings['dowjones'],
    }

    risks = generate_risks()

    opportunities = generate_opportunities()

    social = generate_social()

    governance = generate_governance()

    env_data = {
        'emission_management': generate_env_emission(),
        'resources_management': generate_env_resources(),
        'waste_management': generate_env_waste(),
    }

    presentation = Presentation('template.pptx')
    slide = presentation.slides[0]

    fill_logo(slide, logo_path)

    for shape in slide.shapes:
        if shape.name == 'GeneralInfo':
            fill_data(general_info, shape)
        if shape.name == 'Ratings':
            fill_data(ratings_info, shape)
        if shape.name == 'Risks':
            fill_bullet_list_data(risks, shape)
        if shape.name == 'Opportunities':
            fill_bullet_list_data(opportunities, shape)
        if shape.name == 'Social':
            fill_bullet_list_data(social, shape)
        if shape.name == 'Governance':
            fill_bullet_list_data(governance, shape)
        if shape.name == 'Enviromental':
            fill_env_data(env_data, shape.shapes)

    presentation.save(pptx_path)

    return ratings_info

